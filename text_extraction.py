import os
from PyPDF2 import PdfReader
from pptx import Presentation

def extract_text_from_pdf(pdf_file):
    text = ""
    with open(pdf_file, "rb") as f:
        reader = PdfReader(f)
        num_pages = len(reader.pages)
        for page_num in range(num_pages):
            page = reader.pages[page_num]
            text += page.extract_text()
    return text

def extract_text_from_ppt(ppt_file):
    text = ""
    presentation = Presentation(ppt_file)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text
    return text

def extract_text_from_folder(folder_path):
    extracted_text = {}
    for filename in os.listdir(folder_path):
        file_path = os.path.join(folder_path, filename)
        if filename.endswith(".pdf"):
            extracted_text[filename] = extract_text_from_pdf(file_path)
        elif filename.endswith(".pptx"):
            extracted_text[filename] = extract_text_from_ppt(file_path)
        # Add more conditions for other file types as needed
    return extracted_text
